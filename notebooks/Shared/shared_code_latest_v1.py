# Databricks notebook source
# -*- coding: utf-8 -*-
"""
Created on Mon Feb  3 16:48:39 2020

@author: 809917
"""
#**************************************************
#importing required packages
#**************************************************
import glob
import pytesseract as pyt
import logging
import datetime
import configparser
import shutil
import fitz
import os
from PIL import Image, ImageFilter
from wand.image import Image as wimage
import PIL
import docx
import pptx
import pandas as pd
import PyPDF2 
import openpyxl
import csv
import re
import nltk
from nltk import ngrams
from outlook_msg import Message
import json
import pyodbc
import camelot
import numpy as np

config = configparser.ConfigParser()
#This configuration path should be configured in Blob storage
config.read("/dbfs/mnt/momentive-configuration/config-file.ini")


#Loging environment setup
logger = logging.getLogger('momentive')
logger.setLevel(logging.DEBUG)
fh = logging.FileHandler("shared_main_code.log", 'w')
fh.setLevel(logging.DEBUG)
ch = logging.FileHandler("shared_main_code_error.log", 'w')
ch.setLevel(logging.ERROR)
formatter =logging.Formatter(fmt = '%(asctime)s %(levelname)-8s %(message)s', datefmt='%Y-%m-%d %H:%M:%S')
fh.setFormatter(formatter)
ch.setFormatter(formatter)
logger.addHandler(ch)
logger.addHandler(fh)
text_folder_list = []
image_folder_list = []
record_folder_list = []
invalid_folder_list =[]
#****************************************************
#function name: path_exists
#Objective: To empty or create the folders
#****************************************************
def path_exists(file_path):
  try:
    
    logger.info("Executing path_exists function")
    dbutils.fs.rm(file_path.replace("/dbfs",""),True)
    dbutils.fs.mkdirs(file_path.replace("/dbfs","dbfs:"))
    print('path_exists')
  except Exception as e:
    logger.error("Error in path_exists function : ",exc_info=True)
    
#**********************************************************************
#function name: copy_files
#Objective: To copy files from one folder to another
#input parameter:
#file_list: will hold all the files to moved from target in a list 
#staging_pdf_type: will hold the destination folder path
#called by: sharepoint_native_scanned_pdf_split
#**********************************************************************
def copy_files(file_list, staging_pdf_type):
    try:
        logger.info("Executing copy_files function")
        count=0
        path_exists(staging_pdf_type)
        for file in file_list:
          try:
              file=file.replace("/dbfs","dbfs:")
              file_loc = staging_pdf_type.replace("/dbfs","dbfs:")
              dbutils.fs.cp(file, file_loc)
              logger.info(file + ' copied to ' + staging_pdf_type)
              count+=1
          except Exception as e:
              logger.error("Error while copying ",file)
              logger.error("Error in copy_files function:iteration",exc_info=True)
        logger.info("Number of files copied to "+ staging_pdf_type+" : "+str(count))
    except Exception as e:
        logger.error("Error in copy_files function",exc_info=True)
        
#*******************************************************************************************************************************    
#function name: extract_doc_text
#Objective: To convert documents to text files 
#input parameter:
#all_files: will hold all all-text folder path where extracted text files to be stored  
#staging_path: will hold the staging doc folder path of respective sources
#doc_file_list: will hold all the docx files in a list
#source_type: will hold the respective source type name
#file_processing_info: will hold the query to update extracted text file path  in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database 
#Usage: common code is written which extract text from a document files and store it in a text file on the respective sources
#called by : external_folder_structure_process
#*******************************************************************************************************************************
def extract_doc_text(staging_path,doc_file_list,source_type,all_files,file_processing_info,sql_conn,cursor):
    for files in doc_file_list:
      try:
        doc = docx.Document(files)
        full_text = []
        for text in doc.paragraphs:
            full_text.append(text.text)
            text = '\n'.join(full_text)
        basenames=files.split('/')     
        file_name =  basenames[-1].rsplit('.',1)[0]
        basenames= all_files+basenames[-1].rsplit('.',1)[0]
        text_name = basenames.replace("/dbfs","dbfs:") + '.txt'            
        dbutils.fs.put(text_name,text,True)
        file_path = text_name.replace("dbfs:","/dbfs")
        #**************************************************************************************************************
        #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
        #update_operation
        #***************************************************************************************************************
        file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, {})".format(source_type, file_name, 
        'Document','.docx', staging_path.replace('//','/'), file_path.replace('//','/'), 1,0,'GETDATE()','GETDATE()')
        update_operation(file_processing_info_query,sql_conn,cursor)
        logger.error('{}  extract_doc_text sucessfully'.format(files))
      except Exception as e:
        #**************************************************************************************************************
        #Creation of insert query for the extracted invalid file path to the file_processing_info table and executed using
        #update_operation
        #***************************************************************************************************************
        file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} )".format(source_type, file_name, 
        'Document','.docx', staging_path.replace('//','/'), 'null', 0,0,'null','null')
        update_operation(file_processing_info_query,sql_conn,cursor)
        logger.error('Error in extract_doc_text while processing {}'.format(files))
        



#********************************************************************************************
#This functionality extract text from a powerpoint files and store it in a text file.
#********************************************************************************************
def extract_pptx_text():
    try:
        ppt_nfiles = glob.glob(config.get('path', 'ppt_files') + '*.pptx')
        for files in ppt_nfiles:
            ppt = pptx.Presentation(files)
        
            text_runs = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
                            
            ppt_basenames = os.path.splitext(files)[0].replace('\\', '/').split('/')[-1]
            ppt_name = ppt_basenames + '.txt'
            
            
            f = open(config.get('path', 'extract_path') + ppt_name, 'a+', encoding='utf-8')
            f.write(text_runs)
            f.close() 
    except Exception as e:
        logger.error(e)

#*****************************************************************************************
#function name: sharepoint_native_scanned_pdf_split
#Ojective: Differentiating pdf into Native and Scanned
#staging_folder_path: will hold the staging pdf folder path of respective sources
#raw_pdf_files: will hold all the pdf files in a list
#Usage: Common code is written to split the raw pdf into native and scanned
#called by : external_folder_structure_process 
#******************************************************************************************
def sharepoint_native_scanned_pdf_split(staging_folder_path,raw_pdf_files):
    try:
      count = 0
      scan_files=[]
      native_files=[]
      logger.info("Executing sharepoint_native_scanned_pdf_split function")
      for files in raw_pdf_files:
          try:
              pdfFileObj = open(files, 'rb') 
              #********************************
              # creating a pdf reader object 
              #********************************
              pdfReader = PyPDF2.PdfFileReader(pdfFileObj)  
              if pdfReader.numPages>1:
                  pageObj = pdfReader.getPage(1)
              else:
                  pageObj = pdfReader.getPage(0) 

              if pageObj['/Resources'].get('/XObject') is not None:
                  scan_files.append(files)

              else:
                  native_files.append(files)
              pdfFileObj.close()
          except :
              scan_files.append(files)
              pdfFileObj.close()
      #**********************************************************************************************
      #scan_files: will hold all scanned pdf file path found in the staging folder 
      #copy_files: function will help to copy sacnned pdf from raw folder the scanned-files folder 
      #**********************************************************************************************
      if len(scan_files)>0:        
        logger.info('{} Number of scanned pdf files detected in: {}'.format(len(scan_files),staging_folder_path))
        staging_pdf_scanned = staging_folder_path.rsplit('/',2)[0] + '/scanned-files/'
        copy_files(scan_files, staging_pdf_scanned)
      else:
        staging_pdf_scanned = None
        logger.info('{} Number of scanned pdf files detected in: {}'.format(len(scan_files),staging_folder_path))
        
      #******************************************************************************
      #native_files: will hold all native pdf file path found in the staging folder
      #copy_files: function will help to copy native pdf from raw folder the native-files folder 
      #******************************************************************************  
      if len(native_files)>0:
        logger.info('{} Number of native pdf files detected in: {}'.format(len(native_files),staging_folder_path))
        staging_pdf_native = staging_folder_path.rsplit('/',2)[0] + '/native-files/'
        copy_files(native_files, staging_pdf_native)
      else:
        staging_pdf_native =None
        logger.info('{} Number of native pdf files detected in: {}'.format(len(native_files),staging_folder_path))
      
      return staging_pdf_native, staging_pdf_scanned
    except Exception as e:
      logger.error("Error in sharepoint_native_scanned_pdf_split",exc_info=True)
      
#****************************************************************************************
#function name: intialize_temp_files
#Ojective: Temp folder for temporary execution
#Usage: TO create temporary folders for storing images while converting PDF to Image
#****************************************************************************************
def intialize_temp_files(temp_path):
    try:
        count=0
        logger.info("Executing intialize_temp_files function")
        temp = glob.glob(temp_path + '*.*')  
        if len(temp)==0:
            pass
        else:
            for i in temp:
              i = i.replace("/dbfs","dbfs:")
              dbutils.fs.rm(i)
              count+=1
        logger.info("Number of files got deleted from temporary folder : "+str(count))
    except Exception as e:
        logger.error("Error in initializing temp files function",exc_info=True)

#************************************************************************************************************************************** 
#Function name: native_pdf_extract_text
#Ojective : native pdf files  into text files    
#input parameters:
#native_path : will hold native pdf folder path of respective sources
#all_files: will hold all all-text folder path where extracted text files to be stored  
#staging_path: will hold the staging folder path of respective sources
#source_type: will hold the respective source type name
#file_processing_info: will hold the query to update extracted text file path  in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database  
#Usage: To extract text from scanned pdf files and stores text in a text file on the respective sourcs and insert the text file path
#into file_processing_info table
#called by : external_folder_structure_process
#**************************************************************************************************************************************   
def native_pdf_extract_text(native_path,all_files,staging_path,source_type,file_processing_info,sql_conn,cursor):
    try:
        count=0
        logger.info("Executing native_pdf_extract_text function")
        native_files = glob.glob(native_path + '*.*')  
        logger.info("{} Number of native pdf files from folder {}".format(len(native_files),native_path))
        for files in native_files:
          try:
            text=''
            pdf_file = fitz.open(files)
            n_pages = pdf_file.pageCount
            for n in range(n_pages):
                page = pdf_file.loadPage(n)
                text = text + page.getText()
            basenames=files.split('/')     
            file_name =  basenames[-1].rsplit('.',1)[0]
            basenames= all_files+basenames[-1].rsplit('.',1)[0]
            text_name = basenames.replace("/dbfs","dbfs:") + '.txt'            
            dbutils.fs.put(text_name,text,True)
            file_path = text_name.replace("dbfs:","/dbfs")
            #**************************************************************************************************************
            #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
            #update_operation
            #***************************************************************************************************************
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} )".format(source_type, file_name,      
            'PDF','.pdf', staging_path.replace('//','/'), file_path.replace('//','/'), 1,0,'GETDATE()','GETDATE()')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.info("Successfully extracted {} and updated the file_processing_info table".format(file_name))
            count+=1
          except Exception as e:
          #**************************************************************************************************************
          #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
          #update_operation
          #***************************************************************************************************************
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} )".format(source_type, file_name,     
            'PDF','.pdf', staging_path.replace('//','/'), 'Null', 0,0,'null','null')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.error("Error in native_pdf_extract_text function : iteraion",exc_info=True)
            logger.error("Error while extracting text from native file : {}".format(file_name))
        logger.info("Number of native pdf files got converted into text files successfully : "+str(count))
    except Exception as e:
        logger.error("Error in native_pdf_extract_text function",exc_info=True)
        
#************************************************************************************************************************************** 
#Function name: scanned_pdf_extract_text
#Ojective : Scanned pdf files  into text files  
#input parameters:
#scanned_path : will hold scanned pdf folder path of respective sources
#all_files: will hold all all-text folder path where extracted text files to be stored  
#staging_path: will hold the staging folder path of respective sources
#source_type: will hold the respective source type name
#file_processing_info: will hold the query to update extracted text file path  in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#Usage: To extract text from scanned pdf files and stores text in a text file on the respective sourcs and insert the text file path
#into file_processing_info table
#called by : external_folder_structure_process
#**************************************************************************************************************************************       
def scanned_pdf_extract_text(scanned_path,all_files,staging_path,source_type,file_processing_info,sql_conn,cursor):
    try: 
        logger.info("Executing scanned_pdf_extract_text function")
        scanned_files = glob.glob(scanned_path + '*.pdf')
        logger.info("Number of scanned pdf files from folder "+scanned_path+" : "+str(len(scanned_files)))
        count=0
        temp = all_files.rsplit('/',2)[0] + '/temp/'
        #print(temp)
        for files in scanned_files:
          try:
            #intialize_temp_files(path, temp)
            #intialize_temp_files was replace by temp function
            path_exists(temp)
            #calling pdf to image conversion function
            pdf_to_image_converison(files,temp)
            image_files = glob.glob(temp + '*')
            text_extract = ''
            for j in range(len(image_files)):
                im = PIL.Image.open(image_files[j])
                if im.mode=='P':
                    im = im.convert(palette=0)
                im1 = im.filter(ImageFilter.EDGE_ENHANCE_MORE)                                    
                config1 = (' --psm 6')
                text_val = pyt.image_to_string(im1, config=config1)          
                text_extract = text_extract + text_val
            basenames=files.split('/')       
            file_name =  basenames[-1].rsplit('.',1)[0]
            basenames=all_files+(basenames[-1].rsplit('.',1))[0]
            text_name = basenames.replace("/dbfs","dbfs:") + '.txt'            
            dbutils.fs.put(text_name,text_extract,True)
            file_path = text_name.replace("dbfs:","/dbfs")
            #************************************************************************************************************************
            #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
            #update_operation
            #************************************************************************************************************************
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}' ,{}, {})".format(source_type, file_name,     
            'PDF','.pdf',staging_path.replace('//','/'), file_path.replace('//','/'), 1,0,'GETDATE()','GETDATE()')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.info("Successfully extracted {} and updated the file_processing_info table".format(file_name))
            count+=1
          except Exception as e:
          #************************************************************************************************************************
          #Creation of insert query for the extracted invalid file path to the file_processing_info table and executed using
          #update_operation
          #************************************************************************************************************************
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}',{}, {} )".format(source_type, file_name,     
            'PDF','.pdf', staging_path.replace('//','/'), 'Null', 0,0,'null','null')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.error("Error in scanned_pdf_extract_text function : iteration",exc_info=True)
            logger.error("Error while extracting text from scanned file : {}".format(file_name))
        logger.info("Number of scanned pdf files got converted into text files successfully : "+str(count))
    except Exception as e:
        logger.error("Error in scanned_pdf_extract_text function",exc_info=True)
        
        
        
def eml_text(mount_path, eml_path):
  try:
    logger.info("Executing eml_text function")
    path = config.get(mount_path, eml_path)
    files = glob.glob(path + '*.eml')
    logger.info("Number of email files from folder "+ path +" : "+str(len(files)))    
    for file in files:
      try:
        content = open(file, 'r').read()
        msg = email.message_from_string(content) 
        text = msg.get_payload()[0]
        name = file.split('/')    
        all_file_path = config.get(mount_path,'mnt_mpm2019_all_files')
        if not os.path.exists(all_file_path):
           path_exists(all_file_path)
        name = all_file_path +(name[-1].split('.'))[0]        
        eml_data = name + '.txt'       
        with open(eml_data, 'wb') as f:
          f.write(text.get_payload(decode=True))
      except Exception as e:
        logger.error("Error in eml_text iteration function :",exc_fino=True)
        logger.error("Error while extracting text from email:",file)
  except Exception as e:
    logger.error("Error in eml_text function",exc_info=True)
          
          
def eml_attachment(mount_path, eml_path):
  try:
    logger.info("Executing eml_attachment function")
    path = config.get(mount_path, eml_path)
    files = glob.glob(path + '*.eml')
    logger.info("Number of email files from folder "+path+" : "+str(len(files)))
    for file in files:
      try:
        content = open(file, 'r').read()
        msg = email.message_from_string(content)             
        attachment = msg.get_payload()[1]
        filename = attachment.get_filename()  
        if filename.endswith('.pdf'):
          file_copy_loc = config.get(mount_path, 'mnt_mpm2019_pdf_files')
          f = open(file_copy_loc + filename, 'wb')
          f.write(attachment.get_payload(decode=True))
          f.close()
        else:
          logger.info('{} attachment from outlook not in PDF format so we are not processing'.format(filename))
      except Exception as e:
        logger.error("Error in eml_attachment iteration function :",exc_fino=True)
        logger.error("Error while extracting attachment from email:",file)
  except Exception as e:
    logger.error("Error in eml_attachment function",exc_info=True)


#************************************************************************************************************************************** 
#Function name: outlook_attachment
#Ojective : To fetch the attachments from outlook message 
#input parameters:
#msg_list : will hold outlook mesaage file folder path of respective sources
#staging_path: will hold the staging folder path of respective sources
#raw_files: its a list which will store the raw file path location 
#raw_format: its a list which will store the raw file format
#Usage: common code is written to fetch attachments from outlook message and write into the respective staging pdf raw path 
#called by : external_folder_structure_process
#**************************************************************************************************************************************     
def outlook_attachment(msg_list,staging_path_pdf,raw_files,raw_format):
  try:
    logger.info("Executing outlook_attachment function")   
    logger.info("Number of outlook files from folder : "+str(len(msg_list)))
    pdf_list_outlook = []
    for file in msg_list:
      try:
        with open(file) as msg_file:
            msg = Message(msg_file)
            attach = msg.attachments  
        for i in attach:
          with i.open() as attachment_fp:
              if i.filename.endswith('.pdf'):
                  file_copy_loc = staging_path_pdf
                  pdf_list_outlook.append(file_copy_loc + i.filename)
                  raw_files.append(file_copy_loc + i.filename)
                  raw_format.append('.pdf')
                  logger.info('{} file found in outlook_attachment'.format(i.filename))
                  with open(file_copy_loc + i.filename, 'wb') as my_data: ## write to temporary pdf file
                    my_data.write(attachment_fp.read())
                  logger.info('{} sucessfully written in the path {}'.format(i.filename,file_copy_loc))
              else:
                logger.info('{} attachment from outlook not in PDF format so we are not processing'.format(i.filename))
               
      except Exception as e:
        logger.error("Error in outlook_attachment iteration function :",exc_fino=True)
        logger.error("Error while extracting attachment from outlook:",file)
    logger.info('{} found in msg_list'.format(len(pdf_list_outlook)))
    return pdf_list_outlook 
  except Exception as e:
    logger.error("Error in outlook_attachment function",exc_info=True)

                    
def sharepoint_pdf_files(mpm_staging_pdf_raw):
  try:
    logger.info("Executing sharepoint_pdf_files function")
    raw_pdf_files = glob.glob(mpm_staging_pdf_raw +"*")
   # print('raw_pdf_files', raw_pdf_files)
    if len(raw_pdf_files)>0:
      logger.info("Number of raw pdf files : "+str(len(raw_pdf_files)))
    else:
      logger.error(mpm_staging_pdf_raw+" is empty in sharepoint_pdf_files fucntion")
    native_loc, scanned_loc = sharepoint_native_scanned_pdf_split(mpm_staging_pdf_raw,raw_pdf_files)
    print(native_loc)
    all_files=native_loc.split('staging',1)[0] + 'processed-files/all-text-files/'       
    if native_loc != None:
      native_pdf_extract_text(native_loc,all_files)
    if scanned_loc != None:
      scanned_pdf_extract_text(scanned_loc,all_files)
  except Exception as e:
    logger.error("Error in sharepoint_pdf_files",exc_info=True)
    

#To convert excel file to csv
def excel2csv(path, pathloc, sheetname):
  try:
    logger.info("Executing excel2csv function")
    wb = openpyxl.load_workbook(config.get(path, pathloc))
    sh = wb[sheetname]
    head, tail = os.path.split(config.get(path, pathloc))
    filename = config.get(path, pathloc).split('/')[-1].split('.')[0]
    file = head + '/' + filename + '.csv'
    with open(file, 'w', encoding="utf-8") as f:
        c = csv.writer(f)
        for r in sh.rows:
            c.writerow([cell.value for cell in r])
    return file
  except Exception as e:
    logger.error("Error in excel2csv function",exc_info=True)

#Function to extract and validate tox_list_data 
def tox_list_data(filepath,path,toxic_path):
  try:
    logger.info("Executing tox_list_data function")
    global final_tox_data
    toxic_column=config.get(path,toxic_path).split(",")
    tox_data = pd.read_csv(filepath)
    tox_data.dropna(how='all', axis=0, inplace=True)
    tox_valid_data = tox_data.loc[:,toxic_column]
    final_tox_data = tox_valid_data[~tox_valid_data['Product-Commercial Name'].isnull()]
  except Exception as e:
    logger.error("Error in tox_list_data function",exc_info=True)
  
 #To extract final report date from toxicology files
def toxicology(path,raw_loc,valid_loc,native_loc,scan_loc,ext_loc,temp,mat_loc,elast_loc, toxic):
  try:
      logger.info("Executing toxicology function")
      s_path=config.get(path,raw_loc)
      t_path =config.get(path,valid_loc)
      #to check the t_path if exists
      path_exists(t_path)
      valid_dir=config.get(path,ext_loc)
      ela=config.get(toxic,elast_loc)
      material=config.get(path,mat_loc)
      dir_list=os.listdir(s_path)
      df_mat=pd.read_csv(material,encoding="ISO-8859-1")
      df_ela=pd.read_csv(ela,encoding="ISO-8859-1")
      df_ela_edit=[]
      final_list=[]
      df_ela=df_ela.dropna(subset=['Description'])
      material=["lsr","silsoft"]
      valid_column=["BDT","Nam Prod "]
      headers=["Test Article Number","Study","Description","value_matched_with_ela","column_matched_with_ela","file_matching_value","file_name"]
      df_ela_material=df_ela[~df_ela["Description"].str.contains(material[0],case=False)]
      final_list=pd.DataFrame()
      found_info_list=[]
      for ind in range(len(df_mat)):
        row_list=df_mat.loc[ind].tolist()
        df_each_column=pd.DataFrame()
        final_set_ela=[]
        for n,cl in enumerate(df_mat.columns):           
            try:
                #validate using only BDT and name prod products
                if cl in valid_column:
                  org_value=str(df_mat.loc[ind,cl])
                  value=org_value.strip().lower()
                  if value!='' and value !="nan" and len(value)>0:
                      if value.isdigit():
                        pass
#                           rgx = re.compile(r'(\D{}\D)'.format(value),re.I)
#                           df_ela_edit=df_ela_material[df_ela_material["Description"].str.contains(rgx)]
                      else:
                          rgx = re.compile(r'(([^a-zA-Z]|^){}[^a-zA-Z])'.format(value),re.I)
                          df_ela_edit=df_ela[df_ela["Description"].str.contains(rgx)]
                      if len(df_ela_edit)>0:
                          df_ela_edit["value_Matched_with_ela"]=org_value
                          df_ela_edit["column_Matched_with_ela"]=cl
                          df_each_column=pd.concat([df_each_column,df_ela_edit])                 
            except Exception as e:
                logger.error("Error in toxicology : material validation function",exc_info=True)
                logger.error("Error in validating material : ",str(n))
        if len(df_each_column)>0:   
          df_each_column.drop_duplicates(inplace=True)
          final_list=pd.concat([final_list,df_each_column])
      final_list.drop_duplicates(inplace=True)
      final_list.reset_index(drop=True,inplace=True)
      final_list=final_list.dropna(subset=['Test Article Number'])
      toxic_table=[]
      file_count=0
      #Copying valid toxic files to respective folder                    
      for ind in final_list.index:
          try:
            ela_number=final_list.loc[ind,"Test Article Number"]
            product_name=final_list.loc[ind,"value_Matched_with_ela"]
            final_row_list=final_list.iloc[ind].tolist()
#           check if there is existence
            for file in dir_list:
              try:         
                file_name=file
                if ela_number.lower() in file_name.lower():
                    file_count+=1
                    src=s_path+file_name   
                    dbutils.fs.cp(src.replace('/dbfs','dbfs:'), t_path.replace('/dbfs','dbfs:'))
                    toxic_table.append(final_row_list+["matched with ela number",file_name])
                if product_name.lower() in file_name.lower():
                    file_count+=1
                    src=s_path+file_name               
                    dbutils.fs.cp(src.replace('/dbfs','dbfs:'), t_path.replace('/dbfs','dbfs:'))
                    toxic_table.append(final_row_list+["matched with material number",file_name])
              except Exception as e:
                logger.error("Error in toxicology : file validation function iteration",exc_info=True)
                logger.error("Error in copying valid file",file) 
          except Exception as e:
              logger.error("Error in toxicology : file validation function",exc_info=True)
      logger.info("Toxic matched files count",str(file_count))
      valid_info=pd.DataFrame(toxic_table)
      valid_info.columns=headers
      valid_info.drop_duplicates(inplace=True)
      #calling pdf file extraction function
#       sharepoint_pdf_files(path, valid_loc, native_loc, scan_loc, temp,ext_loc)
      valid_date_file=[]
      valid_file=os.listdir(valid_dir)
      print(len(valid_file))
      for f in valid_file:
        try:
            fpath=os.path.join(valid_dir,f)    
            captured_date=''
            with open(fpath,encoding='utf-8') as file:
                data=' '.join(file)
                tokenize = nltk.word_tokenize(data)
                rgx_pattern=["\d{1,2}\s*\/\d{1,2}\s*\/\d{4}","[a-zA-Z]*\s*\d{1,2}\s*,\s*\d{4}","\d{1,2}(\s*|\-)\w+(\s*|\-)\d{4}"]
                str_pattern=["final issue date","final report date","issue date","report date","completion date","final report"]
                n = 6
                sixgrams = ngrams(tokenize, n)
                for grams in sixgrams:
                    form_text=' '.join(grams)
                    print(form_text)
                    result=re.search("|".join(rgx_pattern), form_text)
                    if(result):
                        for i in str_pattern:
                            if i.lower() in form_text.lower():
                                captured_date=str(result.group().strip())
                                print("form_text",form_text)
                                print(f," - ",captured_date)
                                break
                        if captured_date!='' and len(captured_date)>0:
                            valid_date_file.append([f.replace("txt","pdf"),captured_date])
                            captured_date=''
                            break
        except Exception as e:
            logger.error("Error in toxicology : final date extraction function",exc_info=True)
            logger.error("Error while extracting date from ",f)
          
      #   Final completion date maping
      if len(valid_date_file)>0:
        df_date=pd.DataFrame(valid_date_file)
        df_date.columns=["file_name","final_date"]
        final_df=pd.merge(valid_info, df_date, on='file_name',how='left')
        final_df.to_csv("/dbfstoxic-files/final_date.csv",index=False)
        logger.info("Final key value mapping of final report date of toxicology module has been saved as csv file successfully at /dbfstoxic-files/final_date.csv")
  except Exception as e:
    logger.error("Error in toxicology",exc_info=True)


#*****************************************************************************************************************************
#function name: pdf_to_image_converison
#Objectiv: To convert pdf to image
#input Parameter:
#files: will hold the pdf path which need to be converted into image 
#Ouput parameter:
#target: will hold the ouptut path where converted images will get stored 
#Usage: Common code is written to convert all the pages in the pdf to image in temporary location for tesseract processing
#called by: chemical_structure
#*****************************************************************************************************************************
def pdf_to_image_converison(files,target):
  try:
    logger.info("Executing pdf_to_image_converison function")
    destination=target
    if not os.path.exists(destination):
      os.mkdir(destination)
    with wimage(filename=files, resolution=300) as img:
       img.units = 'pixelsperinch'
       img.compression_quality = 99 
       img.save(filename = destination + '1.png')  
    logger.info("PDF file "+files+" has been converted into image file successfully")
  except Exception as e:
    logger.error("Error in pdf_to_image_converison",exc_info=True)
    logger.error("Error in image file",files)

# To convert image to text with coordinates
def image_to_data_conversion(opened_image,file):
  try: 
    logger.info("Executing image_to_data_conversion function")
    txt=pyt.image_to_data(opened_image)
    txt_read=txt.split('\n')
    coordinates=[]
    last_y1=0
    sentence=''
    word_cords=[]
    line_cords=[]
    line_y1=0
    line_x1=0    
    for i in range(1,len(txt_read)):
      try:
          cords_str=str(txt_read[i]).split('\t')
      #    print(cords_str)
          json={}
          cords={}
          text=cords_str[-1].strip()
          if len(cords_str)>10 and len(text)>0:
              word=cords_str[-1]
              json['text']=word
              y1=int(cords_str[7])
              cords["x1"]=cords_str[6]
              cords["y1"]=cords_str[7]
              cords["x2"]=int(cords_str[8])+int(cords_str[6])
              cords["y2"]=int(cords_str[9])+int(cords_str[7])
              json["coordinates"]=cords
              y_dif=(last_y1-y1)
              if y_dif <0:
                  y_dif=-1*y_dif
              if y_dif <=20  and last_y1>0:
                  sentence+=word+" "
                  word_cords.append(json)
              else:
                  if len(word_cords)>0:
#                       print(sentence)
                      line_json={"text":sentence.strip(),
                                 "coordinates":{"x1":line_x1,"y1":line_y1,
                                                "x2":word_cords[-1]["coordinates"]["x2"],
                                                "y2":word_cords[0]["coordinates"]["y2"]}}
                      line_cords.append(line_json)
                      final={"line_cords":line_cords,
                             "word_cords":word_cords}
                      word_cords=[]
                      line_cords=[]
                      coordinates.append(final)
                  line_x1=int(cords_str[6])
                  line_y1=int(cords_str[7])           
                  sentence=''
                  sentence+=word+" "
                  last_y1=y1
                  word_cords.append(json)
      except Exception as e:
        logger.error("Error in image to data conversion: inner iteration",exc_info=True)
    line_json={"text":sentence, 
               "coordinates":{"x1":line_x1,
                              "y1":line_y1,
                              "x2":word_cords[-1]["coordinates"]["x2"],
                              "y2":word_cords[-1]["coordinates"]["y2"]}}
    line_cords.append(line_json)
    final={"line_cords":line_cords,
           "word_cords":word_cords}
    coordinates.append(final)
    logger.info("Text data with their coordinates has been extracted successfully from image file "+file)
    return coordinates
  except Exception as e:
    logger.error("Error in image_to_data_conversion",exc_info=True)
    logger.error("Error in image file",file)

#***************************************************************************************************************************************
#function name: chemical_structure
#Objectiv: To extract chemical structure from file
#input parameters:
#unstruct_data_key_info: will hold all the data except key_value extract data like(product_type, category, product)
#raw_df: will hold all the staging file path in dataframe which helps to move file to processed folder
#data_extract: will hold the key-value data 
#Usage: common code is written which extracts chemical structure for the identified product in the files based on the coordinates data #produced by tesseract ocr and update unstruct_data_key_info dataframe       
#called by: key_value_extract
#***************************************************************************************************************************************
def image_data_extract(file):
#def chemical_structure(unstruct_data_key_info,index,data_extract,raw_df):
#def chemical_structure(chemical_path,input_path,image_path,converted_image,inscope_structure):
  try:        
        file=file.replace("dbfs:","/dbfs").strip()      
        #file_name =  file.split('/')[-1][:-4]
        logger.info("Executing image_Key_extract function")
        check_path=True
        temp_path = file.rsplit('staging',1)[0]  + 'temp/'
        intialize_temp_files(temp_path)
        #Conveting pdf to image file
        logger.info("Calling pdf to image conversion function")
        pdf_to_image_converison(file,temp_path)
        target= temp_path
        target_list= glob.glob(target+'*.*')
        return target_list
  except Exception as e:
     pass
      

def image_to_cordinates(sql_conn,cursor,img_path,product_type_list,product_list,file_loc,category,file_name,img_count,unstructure_processed_data_query):
  block_json = {}
  json_check = []
  print(product_list)
  #global img_count
  try:
      im = Image.open(img_path) 
      width, height=im.size
      #Converting image to text with coordinates as json doc
      logger.info("Calling image to data conversion function")
      coordinates=image_to_data_conversion(im,img_path)
      crop_json={}
      crop_details=[]
      first_count=0
      expected_left=0
      diff_x1=0
      for item in coordinates:
        try:
          line=item['line_cords']
          for ele in line:
              text=ele["text"].strip()
             # print(text)
              x1=int(ele["coordinates"]["x1"])
              y1=int(ele["coordinates"]["y1"])
              x2=int(ele["coordinates"]["x2"])
              y2=int(ele["coordinates"]["y2"])
             # print(text.lower().strip())
             # rgx_img = re.compile('(.|\n|\t\r)*{}(.|\n|\t\r)*',re.I)
              match_f = None
              for prod_name in product_list:
                rgx_img = re.compile('{}'.format(prod_name.replace('*','\*')),re.I)
                for match in re.finditer(rgx_img,text):
                  match_f = prod_name
                  
                  print('print(match_f)',match_f)
              if  match_f is not None:
                  
                  first_count+=1
                  if first_count==1:
                      expected_left=x1
                  left=x1
                  top=y1-5
                  if(crop_json):
                      crop_json["bottom"]=y1-5
                      crop_details.append(crop_json) 
                  crop_json={}
                  crop_json={"left":0,
                            "top":top,
                            "right":width,
                              "name":match_f}

              elif(len(text)>1 and first_count>1):
                  diff_x1=expected_left-x1
                  if diff_x1<0:
                      diff_x1=diff_x1*-1
                  if (diff_x1<=3):
                      crop_json["bottom"]=y1-5
                      crop_details.append(crop_json)
                      crop_json={}                        
              if first_count==1:
                  first_count+=1
        except Exception as e:
          logger.error("Error in extracting image_Key_extract function: coordinates iteration",exc_info=True)
      if(crop_json):
          if (y2-y1)>50:
              crop_json["bottom"]=y2+5
          else:
              crop_json["bottom"]=height
          crop_details.append(crop_json)
      else:
        crop_json={"left":0,
                    "top":0,
                    "right":width,
                    "bottom" : height,
                      "name":"file_name"}
        crop_details.append(crop_json)
#     except Exception as e:
#         logger.error("Error in extracting image_Key_extractfunction: finding coordinates of chemical structure iteration",exc_info=True)
#         logger.error("Error in extracting file :",tfile)
    #logger.info("image_Key_extract image coordinates has been founded")

      if len(crop_details)>0:
        
        for item in crop_details:
            im1 = im.crop((item.get('left'), item.get('top'), item.get('right'), item.get('bottom'))) 
            name=item.get('name')
            name_db = name
            if name == 'file_name':
              img_path = file_loc + 'non-relavent/' 
              img_count = img_count + 1
              name = file_name + '_' + str(img_count)
              product_type = 'null'
              name_db = 'null'
            else:
              prod_index = product_list.index(name)
              product_type = product_type_list[prod_index]
              img_path = file_loc + 'relavent/' 
            if not os.path.exists(img_path) :
              path_exists(img_path)
              check_path=False
            path=img_path+name+".png"
            im1.save(path)
            block_json['file_path'] = path
            data_extract=json.dumps(block_json)
            if product_type == 'null':
              unstructure_processed_data(unstructure_processed_data_query,category,product_type,name_db,data_extract,0,sql_conn,cursor)
            else:
              unstructure_processed_data(unstructure_processed_data_query,category,product_type,name_db,data_extract,1,sql_conn,cursor)
#             logger.info("{} image : ".format(unstruct_data_key_info['category'][index])+path+" saved successfully")
#             block_json['file_path'] = path
#             if unstruct_data_key_info.shape[0] > len(data_extract): 
#               data_extract.append(json.dumps(block_json))                 
#             else:
#               data_extract.append(json.dumps(block_json))
#               unstruct_data_key_info['file']
#               unstruct_data_df['file_path'].append(str_file)  
#               unstruct_data_df['file_type'].append('image')  
#               unstruct_data_df['category'].append(category_list)    
#               unstruct_data_df['product_type']  = product_type_list
#               unstruct_data_df['product']  = product_valid_list 
          # else:
          #     img_path = file.rsplit('/',2)[0].replace("dbfs:","/dbfs") + '/key-data-extract/'
          #     if not os.path.exists(img_path): 
          #       os.mkdir(img_path)
          #       #path_exists(img_path)
          #       check_path=False
          #     path=img_path+ tfile
          #     dbutils.fs.cp(os.path.join(target,tfile).replace("/dbfs","dbfs:"),img_path.replace("/dbfs","dbfs:"))
          #     logger.info("{} image : ".format(unstruct_data_key_info['category'][index]) +path+" saved successfully")
          #     block_json['file_path'] = path                     
          #     data_extract.append(json.dumps(block_json))
      return img_count
  except Exception as e:
    logger.error("Error in extracting chemical structure",exc_info=True)

 

#***************************************************************************************************************************************
#function name: molecular_weight
#Objectiv: To extract molecular weight from file
#input parameters:
#unstruct_data_key_info: will hold all the data except key_value extract data like(product_type, category, product)
#raw_df: will hold all the staging file path in dataframe which helps to move file to processed folder
#data_extract: will hold the key-value data 
#Usage: common code is written which extracts required key value data based on the keywords for mol-weight categories and update the 
#       unstruct_data_key_info dataframe 
#called by: key_value_extract
#*************************************************************************************************************************************** 
def molecular_weight(unstruct_data_key_info,index,data_extract,raw_df):
#def molecular_weight(molecular,file_path,text_path,native,scan,temp):
  try:
    #sending to sharepoint_pdf_files function for text extraction
    logger.info("Executing molecular_weight function")
    logger.info("calling sharepoint_pdf_files funtion for pdf to text extraction")
    file=unstruct_data_key_info['file_path'][index].replace("dbfs:","/dbfs")
    content = open(file, 'r', encoding = 'utf-8').read().lower()
    weight=''
    mono_weight=''
    avg_weight=''
    random=''
    block_json ={}
    with open(file) as fobj:
      for i in fobj:
        try:
            if "molecular weight" in i.lower():  
              relist=["((\d+.\d+)|\d+)\s*((\w+/\w+)|\w+)","((\d+.\d+)|\d+).*"]
              rgx=re.search("|".join(relist),i,re.I)
              if rgx is not None:
                if "monoisotopic" in i.lower():
                  mono_weight=rgx.group().strip()
                  block_json['mol_weight'] = mono_weight
                elif "average" in i.lower():
                  avg_weight=rgx.group().strip()
                  block_json['mol_weight'] = avg_weight
                else:
                  random=rgx.group().strip()
                  block_json['mol_weight'] = random
        except Exception as e:
            logger.error("Error in molecular weight inside function",exc_info=True)
            logger.error("Error in finding molecular weight of file :",file)
    data_extract.append(json.dumps(block_json))
    
  except Exception as e:
    logger.error("Error in molecular weight function",exc_info=True)
    
def registration_status(path,file_path,eu,us,latam):
  try:
#       EU_extract_data=pd.read_csv(excel2csv(path,file_path,'LATAM'))
      file=config.get(path, file_path)
      eu_column=config.get(path,eu)
      us_column=config.get(path,us)
#       print("org",(us_column.split(',')))
      latam_column=config.get(path,latam)
#       print(eu_column)
      EU_extract_data = pd.read_csv(file.replace('.xlsx','_EU.csv'))
      US_extract_data = pd.read_csv(file.replace('.xlsx','_US.csv'), header=None)
      LATAM_extract_data = pd.read_csv(file.replace('.xlsx','_LATAM.csv'))

      EU_data_rel = EU_extract_data.loc[:,eu_column.split(',')]
#       print(US_extract_data.columns)
      US_data_filter = US_extract_data[1:]
#       print("nn",US_data_filter.columns)
      US_data_col = US_data_filter.rename(columns=US_data_filter.iloc[0])
      US_data_col.drop(US_data_col.head(1).index, inplace=True)
#       print("uss",US_data_col.columns)
      US_data_rel = US_data_col.loc[:,us_column.split(",")]
#       print(US_data_rel.columns)
      LATAM_data_rel = LATAM_extract_data.loc[:,latam_column.split(",")]
#       print(LATAM_data_rel.columns)
#       US_products = ['Silwet* 408', 'Silwet* L-77', 'Silwet* 806', 'Y-19334']
#       EU_products = ['Silwet 806', 'Silwet L-77', 'Silwet 408', 'Y-19334']
#       LATAM_products = ['Silwet 806', 'Silwet L-77', 'Silwet 408', 'Y-19334']

#       LATAM_final_df = pd.DataFrame() 
#       for i in range(len(LATAM_products)):
#           l = LATAM_data_rel[LATAM_data_rel.Product==LATAM_products[i]]
#           LATAM_final_df = pd.concat([l, LATAM_final_df])


#       US_final_df = pd.DataFrame()
#       for i in range(len(US_products)):
#           m = US_data_rel[US_data_rel['Product Name']==US_products[i]]
#           US_final_df = pd.concat([m, US_final_df])

#       EU_final_df = pd.DataFrame()
#       for i in range(len(EU_products)):
#           e = EU_data_rel[EU_data_rel['Product']==EU_products[i]]
#           EU_final_df = pd.concat([e, EU_final_df])
  except Exception as e :
    logger.error("Error in registration-status",e)

def writing_invalid_into_file(module_name,path,inuput_path,output_path,file_extension=''):
  try:
    logger.info("Executing writing_invalid_into_file function")
    folder=os.listdir(config.get(path,inuput_path))
    output=config.get(path,output_path)
    folder_length=len(folder)
    
    if file_extension!='':
      total_files="\n".join(folder)
      file_extension="."+file_extension
      replaced_files=total_files.replace(".txt",file_extension)  
#     out_file_name=(output+module_name+".txt").replace("/dbfs","dbfs:")
    out_file_name=output+module_name+".txt"
    final_str="Number of invalid files : {}".format(folder_length)+"\n"+replaced_files
#     dbutils.fs.put(out_file_name,final_str,True)
    with open(out_file_name,"w") as file_object:
      file_object.write(final_str)
  except Exception as e:
    logger.error("Error in writing_invalid_into_file function",e)
#unstructure_processed_data_query,category,'NAM PROD',nam_prod,data_extract,sql_conn,cursor
def unstructure_processed_data(unstructure_processed_data_query,category,product_type,product,data_extract,is_relevant,sql_conn,cursor):
  #print(unstruct_key_value_data_df.columns)
  #for index in unstruct_key_value_data_df.index:
    insert_query = unstructure_processed_data_query.format(category,product_type,product,data_extract,'getdate()','getdate()',is_relevant)
    #print(insert_query)
    update_operation(insert_query,sql_conn,cursor) 

#*********************************************************************************************************************************
#function name: tex_Key_extract
#Ojective: Key value extart for FDA files
#input parameters:
#unstruct_data_key_info: will hold all the data except key_value extract data like(product_type, category, product)
#raw_df: will hold all the staging file path in dataframe which helps to move file to processed folder
#data_extract: will hold the key-value data 
#Usage: common code is written which extracts required key value data based on the keywords for US-FDA categories and update the 
#       unstruct_data_key_info dataframe
#called by: key_value_extract
#*********************************************************************************************************************************** 
def text_Key_extract(file,filter_df,content):
#def text_Key_extract(unstruct_data_key_info,index,data_extract,raw_df,filter_df):
  try: 
    file=file.replace("dbfs:","/dbfs").strip()
    print(file)
    file_name =  file.split('/')[-1].rsplit('.',1)[0]
    #staging_raw_file_path will contain raw file path  to move into the processed folder
    #staging_raw_file_path = raw_df[raw_df['file_name'].str.contains(file_name)]['file_name']
    block_json = {}
    #******************
    #value_extract
    #******************
    for  index_df in filter_df.index:
        start_string_index = None
        end_string_index = None
        rgx = re.compile(r'({})'.format(filter_df['start_key'][index_df]),re.I)
        #******************************
        #checking index of start key
        #******************************
        for match in re.finditer(rgx,content):
            if match.group():
                start_string_index = match.start()
                break

        if  start_string_index is not None:
            #******************************
            #checking index of end key
            #******************************
            if filter_df['end_key'][index_df] != 'date_type':
                rgx = re.compile(r'({})'.format(filter_df['end_key'][index_df]),re.I)  
            else:
              #  \d{1,2}\s*\/\d{1,2}\s*\/\d{4}","[a-zA-Z]*\s*\d{1,2}\s*,\s*\d{4}","\d{1,2}(\s*|\-)\w+(\s*|\-)\d{4}
                rgx ="|".join(["\d{1,2}\s*\/\d{1,2}\s*\/\d{4}","[a-zA-Z]*\s*\d{1,2}\s*,\s*\d{4}","\d{1,2}\s[a-zA-Z-]*\s\d{4}","\d{1,2}\s*\-\s*[a-zA-Z]*\s*\-\s*\d{4}"])
            for match in re.finditer(rgx,content[start_string_index:]):
                if match.group():
                    end_string_index = start_string_index + match.end()  
                    break
                    
        if  start_string_index is not None and end_string_index is not  None:
            text_extract = content[start_string_index:end_string_index].replace('\n',' ')#.replace('\u2014'
        
            #********************************************************************************
            #replacing end_key text in extracted text if it is presnt in start key column
            #********************************************************************************
            if not filter_df[filter_df['start_key'].str.contains(filter_df['end_key'][index_df])].empty:
                find_replace=text_extract.lower().find(filter_df['end_key'][index_df].lower())
                text_extract = text_extract[:find_replace]

            #**********************************************************************************
            #Writing the extracted text in the json based on the field names present in the 
            #***********************************************************************************
            
            if pd.isnull(filter_df['field'][index_df]):
                if 'data' not in block_json.keys():
                    block_json['data'] = text_extract
                else:
                    block_json['data'] = block_json['data'] + ' ' + text_extract
            else:
                if filter_df['field'][index_df] not in block_json.keys():
                    block_json[filter_df['field'][index_df]] = text_extract
                else:
                    block_json[filter_df['field'][index_df]] = block_json[filter_df['field'][index_df]] + ' ' + text_extract
                   
    #****************
    #date extract
    #****************
    rgx_pattern_date =["\d{1,2}\s*\/\d{1,2}\s*\/\d{4}","[a-zA-Z]*\s*\d{1,2}\s*,\s*\d{4}","\d{1,2}\s[a-zA-Z-]*\s\d{4}","\d{1,2}\s*\-\s*[a-zA-Z]*\s*\-\s*\d{4}"] 
    date_result=re.search("|".join(rgx_pattern_date), content)
    folder_path = file.rsplit('/',2)[0] + '/key-data-extract/'   
    if not os.path.exists(folder_path):
      os.mkdir(folder_path)
    text_name =  folder_path + file.split('/')[-1][:-4] + '.json'    
    if date_result is not None:
       block_json['Date'] = date_result.group().strip()
    else:
      block_json['Date'] = 'null'
      
    #********************
    # subject extract:
    #********************
    rgx_pat = r'\nre(:|\s)'
    sub_first = re.finditer(rgx_pat, content)
    for m_string in sub_first:
        sub_first_check_v1 = m_string.start(0)
        break
    sub_first_check_v2 = content.find('eu food contact statement')
    sub_last = content.find('\ndear')
    sub_first_check_v3 = content.rfind('\n',0,sub_last)    
    if sub_last != -1 and sub_first_check_v1 != -1 and sub_last > sub_first_check_v1 :
        subject = content[sub_first_check_v1:sub_last].strip()
    elif sub_last != -1 and sub_first_check_v2 != -1 and sub_last > sub_first_check_v2 :
        subject = content[sub_first_check_v2:sub_last].strip()
    elif sub_last != -1 and sub_first_check_v3 != -1 and sub_last > sub_first_check_v3 :
        subject = content[sub_first_check_v3:sub_last].strip()
    else:
        subject = file.split('/')[-1][:-4]
        block_json['subject'] = subject
    
    #*****************************************
    #Copying files to processed folder
    #*****************************************
#     processed_path = file.rsplit('analytics',1)[0]  + 'Processed/'
#     if not os.path.exists(processed_path.replace("dbfs:","/dbfs")):
#         os.mkdir(processed_path.replace("dbfs:","/dbfs"))
#     dbutils.fs.cp(staging_raw_file_path.values[0].replace("/dbfs","dbfs:").replace('//','/'), processed_path.replace("/dbfs","dbfs:")) 
#     block_json['file_path']  = processed_path.replace("dbfs:","/dbfs") + staging_raw_file_path.values[0].rsplit('/',1)[1]                  
#     data_extract.append(json.dumps(block_json))  
    
#     with open(text_name, "w",encoding='utf8') as file_write:
#       json.dump(block_json,file_write,ensure_ascii=False)
    logger.info("data extract is successful for this {}".format(file,text_name))  
    return block_json
 
  except Exception as e:
    logger.error("Error in text_Key_extract function while processing {}".format(file),exc_info=True)
    
    

#*************************************************************************************************************************
#Function name: key_value_extract 
#Ojective: Key value extract on relevant and non-relevant files
#input parameters:
#unstruct_data_key_info: will hold all the data except key_value extract data
#raw_df: will hold all the staging file path in dataframe which helps to move file to processed folder
#Usage: This functionality does the calling of key-value extract function baed on categories
#output parameters: 
#unstruct_data_key_info: will hold all the unstructure consolidate data
#called by: main function
#*************************************************************************************************************************
def key_value_extract(file,content,category,unstruct_category_key_df,file_type,raw_df):
#def key_value_extract(unstruct_data_key_info,raw_df,unstruct_category_key_df):
 try:    
    #filter_df will contain the start key and  end key for the data to be extracted from the file
    filter_df= unstruct_category_key_df[unstruct_category_key_df['category'].str.contains(category)]
    if not filter_df.empty:    
        if file_type.strip() == 'text':
            data_extract=text_Key_extract(file,filter_df,content)

        elif file_type.strip().strip() == 'image':
            print('image')
            image_data_extract(file,filter_df)    

        elif file_type.strip().strip().strip() == 'records':
            record_Key_extract(unstruct_data_key_info,index,data_extract,raw_df,filter_df)    

#     print(unstruct_data_key_info.shape)
#     print(len(data_extract))
#     unstruct_data_key_info['data_extract'] = data_extract 
    return data_extract
 except Exception as e:
    logger.error("Error in key_value_extract function while processing {}".format(file),exc_info=True)
    
#**************************************************************************************************************************************
#function name: file_validation
#Ojective: Moving the files to relevant and non-relevant folder on the respective sources
#input parameter: 
#file: will hold the absolute file path of all-text folder
#file_valid_type: will hold the valid folder file path 
#file_is_valid_query: will hold the query to update the file_processing_info table for valid(relevant and non-relevant file path) and #                     non-valid file path
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#Usage: This functionality helps to moves the files to relavent and non-relavent folder based on the function call done on the 
#       relavent_file_extract
#       To update the file_processing_info table
#called by: relavent_file_extract function
#************************************************************************************************************************************** 

def valid_files_copy(file,file_valid_type,data_extract):    
  try:
    if not os.path.exists(file_valid_type):
      path_exists(file_valid_type)
    text_name = file_valid_type.replace("dbfs:","/dbfs") + file.split('/')[-1][:-4] + '.txt'
    with open(text_name, "w",encoding='utf8') as file_write:
       json.dump(data_extract,file_write,ensure_ascii=False)
    logger.info('key-data extract of file {} has been written into {}'.format(file,text_name))    
  except Exception as e:
    logger.error('something went wrong in file_validation function',exc_info=True)
    
def relavent_image_extract(file,file_loc,content,product_inscope_df,category,file_is_valid_query,file_unique_list,sql_conn,
                           cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query):  
  global image_folder_list
  
  try: 
      print(image_folder_list)       
      if file_loc.strip() not in image_folder_list :
        image_folder_list.append(file_loc.strip())
        path_exists(file_loc)
      prod_flag =''
      product_list = []
      product_type_list =[]
      data_extract={}
      target_temp=image_data_extract(file)    
      file_name =  file.split('/')[-1].rsplit('.',1)[0]
      img_count = 0
      
      
      #data_extract.append(json.dumps(block_json))  
      #*********************************************************
      #checking the  PROD present in the exrtacted content 
      #********************************************************* 
      for prod_index in product_inscope_df.index:
        if not pd.isnull(product_inscope_df['Type'][prod_index]):
           try:
              prod_rgx = re.search(r'(([^a-zA-Z]|^){}[^a-zA-Z])'.format(re.escape(product_inscope_df['Text'][prod_index].strip())),content,re.I)   
              if(prod_rgx):
                  prod_txt =  product_inscope_df['Text'][prod_index].strip().upper()
                  prod_flag = 's'                
                  logger.info('{} Successfully passed the inscope validation by containing {} {} in the content'.format(file,product_inscope_df['Type'][prod_index].strip(),prod_txt)) 
                  product_type_list.append(product_inscope_df['Type'][prod_index].strip())
                  product_list.append(prod_txt)
           except Exception as e:
                logger.error("Error in relavent_image_extract function inner iteration",exc_info=True) 
              
      for img_path in target_temp:
        img_count = image_to_cordinates(sql_conn,cursor,img_path,product_type_list,product_list,file_loc,category,
                                        file_name,img_count,unstructure_processed_data_query)

               
  except Exception as e:
      logger.error("Error in relavent_image_extract function: outer iteration",exc_info=True)
          
  
#**************************************************************************************************************************************
#function name: relavent_file_extract
#Ojective: TO identify relevant and non-relevant files
#input parameter: 
#file: will hold the absolute file path of all-text folder
#file_loc: will hold the valid folder file path 
#bdt_list: will hold all the BDT data fetched from product_inscope_df
#nam_prod_list: will hold all the NAM PROD data fetched from product_inscope_df
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#category_list: will store the category type in which extracted data falls into(like BDT, NAM PROD)
#product_type_list: will store product type which identified in extracted data
#file_path_list: will store absolute file path of all-text folder
#Usage: To differentiate relevant and non-relavent files based on the product inscope details like(NAM PROD, BDT, SILICONE US-FDA and EU
#file_validation: Moves the file to relevant and non-relevant folder based on the product_inscope
#called by: pattern_match_validation
#**************************************************************************************************************************************   
def relavent_text_extract(file,file_loc,content,product_inscope_df,category,file_is_valid_query,file_unique_list,sql_conn,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query,sil_elast_product_list= None):
  global text_folder_list
  try: 
      file_relavent = file_loc + 'relavent/'
      file_non_relavent = file_loc + 'non-relavent/'
      prod_flag =''      
      sil_prod_flag =''
      data_extract={}
      filter_df= unstruct_category_key_df[unstruct_category_key_df['category'].str.contains(category)]
      if not filter_df.empty:
        data_extract=text_Key_extract(file,filter_df,content)
            

     #*****************************************
     #Copying files to processed folder
     #*****************************************
      file_name =  file.split('/')[-1].rsplit('.',1)[0]
      staging_raw_file_path = raw_df[raw_df['file_name'].str.contains(file_name)]['file_name']
      processed_path = file.rsplit('analytics',1)[0]  + 'Processed/'
      if file_loc.strip() not in text_folder_list :
        text_folder_list.append(file_loc.strip())
        path_exists(file_loc)
      if not os.path.exists(processed_path):        
        path_exists(processed_path)
      dbutils.fs.cp(staging_raw_file_path.values[0].replace("/dbfs","dbfs:").replace('//','/'), processed_path.replace("/dbfs","dbfs:")) 
      data_extract['file_path']  = processed_path.replace("dbfs:","/dbfs") + staging_raw_file_path.values[0].rsplit('/',1)[1]     
      data_extract = json.dumps(data_extract)
      
           
      #data_extract.append(json.dumps(block_json))  
      #*********************************************************
      #checking the NAM PROD present in the exrtacted content 
      #*****************************************************prod_index**** 
      for prod_index in product_inscope_df.index:
        
        if not pd.isnull(product_inscope_df['Type'][prod_index]):
           try:
              if not str(product_inscope_df['Text'][prod_index].strip()).isspace():
                prod_rgx = re.search(r'(([^a-zA-Z]|^){}[^a-zA-Z])'.format(re.escape(product_inscope_df['Text'][prod_index].strip())),content,re.I)   
                if(prod_rgx):
                    print('ll',product_inscope_df['Text'][prod_index].strip())
                    prod_txt =  product_inscope_df['Text'][prod_index].strip().upper()

                    prod_flag = 's'                
                    logger.info('{} Successfully passed the inscope validation by containing {} {} in the content'.format(file,product_inscope_df['Type'][prod_index].strip(),prod_txt))   
                    unstructure_processed_data(unstructure_processed_data_query,category,product_inscope_df['Type'][prod_index].strip(),prod_txt,data_extract,1,sql_conn,cursor)
           except Exception as e:
                logger.error("Error in relavent_file function inner iteration",exc_info=True)  
      
      #*************************************************************************************************************
      #checking the EU and US-FDA present in the exrtacted content if NAM PROD and BDT not in the content
      #*************************************************************************************************************
      if prod_flag != 's'  and sil_elast_product_list != None:            
          try:
            for sil_prod in sil_elast_product_list:
                sil_prod_rgx = re.search(r'(([^a-zA-Z]|^){}[^a-zA-Z])'.format(re.escape(sil_prod)),content,re.I)          
                if(sil_prod_rgx):
                    sil_prod_flag = 's'
                    sil_prod = sil_prod.upper()
                    logger.info('{} Successfully passed the inscope validation by containing silicone elatomer  {} in the content'.format(file,sil_prod))
                    unstructure_processed_data(unstructure_processed_data_query,category,'Silicone elastomer','NAM PROD or BDT',data_extract,1,sql_conn,cursor)
          except Exception as e:
            logger.error("Error in relavent_file function: silicone_elast_prod inner iteration",exc_info=True)
      
      #*************************************************************************************************************
      #Moving the files to Non-relevant folder if NAM PROD, BDT, EU and US-FDA not in the content
      #*************************************************************************************************************
      if prod_flag != 's'   and sil_prod_flag != 's' :         
            logger.info('{} it does not fall under incsope product, So moving this to {}'.format(file,file_non_relavent))
            valid_files_copy(file,file_non_relavent,data_extract) 
            unstructure_processed_data(unstructure_processed_data_query,category,'null','null',data_extract,0,sql_conn,cursor)
            #logger.info('{} copying the files to non relevant {}'.format(file, file_non_relavent))
      else: 
            valid_files_copy(file,file_relavent,data_extract)
            #logger.info('{} copying the files relevant {}'.format(file, file_relavent))
      #text_flag = 's'
  except Exception as e:
      logger.error("Error in relavent_file function: outer iteration",exc_info=True)
          
#******************************************************************************************************************************************
#function name : pattern_match_validation
#Ojective: Pattern match validation on each file extracted from staging path
#input parameters:
#external_processed_files: will have all the extracted file paths from the file_processing_info table
#external_staging_path: will have all the staging folder path from the file_processing_info table for the extracted file
#external_satging_file_format: will have all file formats from the file_processing_info table for the extracted file
#pattern_matching_query: will have select query for the pattern_matching_keys table
#pattern_key_df: will store all pattern matching keywords by passing pattern_matching_query to external_source_data function
#product_inscope_query: will have the select query of product_inscope(Nam prod, BDT,CAS No details)
#product_inscope_df: will hold all product inscope details in dataframe by passing product_inscope_query to external_source_data function
#bdt_list: will hold all the BDT data fetched from product_inscope_df
#nam_prod_list: will hold all the NAM PROD data fetched from product_inscope_df
#silicone_elastomer_product_query: will have select query for the silicone elastomer table which contains 
#                                  valid US-FDA and EU Product data extracted from Silicone elastomer brochure pdf file
#                                  scrapped from momentive website
#silicone_elastomer_product_df: will hold all the US-FDA and EU in dataframe by passing product_inscope_query to external_source_data 
#                               function
#relavent_file_extract: To differentiate relevant and non-relavent files based on the product inscope details
#Usage: This functionality helps to perform 15 pattern match validation on each extracted text files based on all_text path stored in #file_processing_info table
#output parameter:
#unstruct_data_df: will hold all the valid file paths, categories and product keys which will be input for key value extract function
#******************************************************************************************************************************************   
def pattern_match_validation(sql_conn,external_processed_files_df,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query):
  try:
    external_processed_files = external_processed_files_df['blob_all_txt_file_path'].values.tolist()
    external_staging_path = external_processed_files_df['blob_staging_file_path'].values.tolist()
    external_satging_file_name = external_processed_files_df['file_name'].values.tolist()
    external_satging_file_format = external_processed_files_df['file_format'].values.tolist()
    pattern_matching_query = config.get('mount_path','pattern_match')  
    pattern_key_df = external_source_data(sql_conn,pattern_matching_query)  
    pattern_category = list(set(pattern_key_df['pattern_category'].values.tolist()))
    file_is_valid_query = config.get('mount_path', 'file_is_valid')
    product_inscope_query = config.get('mount_path','product_inscope')
    product_inscope_df = external_source_data(sql_conn,product_inscope_query)
    product_inscope_df.replace(r'^\s*$', np.nan, regex=True,inplace=True)
    product_inscope_df = product_inscope_df.fillna('null')
    product_inscope_df = product_inscope_df.fillna('null')
    silicone_elastomer_product_query = config.get('mount_path','silicone_elastomer_product')
    silicone_elastomer_product_df =  external_source_data(sql_conn,silicone_elastomer_product_query)    
    bdt_list = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('MATNBR')]['Text3'].values.tolist()))
    nam_prod_list = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('NAMPROD')]['Text1'].values.tolist()))
    cas_list = list(set(product_inscope_df[product_inscope_df['Type'].str.contains('NUMCAS')]['Text1'].values.tolist()))
    product_type_list= ['NAMPROD'] * len(nam_prod_list) + ['BDT'] * len(bdt_list) +  ['NUMCAS'] * len(cas_list)
    product_valid_list = nam_prod_list + bdt_list + cas_list
    product_inscope_df = pd.DataFrame(columns=['Type', 'Text'])
    product_inscope_df['Type'] = product_type_list
    product_inscope_df['Text'] = product_valid_list
    #product_inscope_df['NUMCAS'] = cas_list
    silicone_elastomer_product_query = config.get('mount_path','silicone_elastomer_product')
    silicone_elastomer_product_df =  external_source_data(sql_conn,silicone_elastomer_product_query)  
    silicone_elastomer_product_df=silicone_elastomer_product_df.rename(columns = {'eu_fda':'EU-FDA','us_fda':'US-FDA'})
    file_unique_list =[]
    category_list = []
    product_valid_list =[]
    product_type_list =[]
    file_path_list= []
    file_type_list = []
    global invalid_folder_list
    unstruct_data_df = pd.DataFrame(columns=['file_path', 'category', 'product_type','product'])
    
    #**********************************************
    #Iterating each files for pattern matching 
    #**********************************************
    for index in range(0,len(external_processed_files)):
      try:
        file=external_processed_files[index].replace("dbfs:","/dbfs")
        content = open(file, 'r', encoding = 'utf-8').read()
        file_valid_flag =''
        #print(content)
        
        #****************************************************************************************************
        #checking pattern_match on each file based on pattern_category,pattern_keys and filter_condition
        #***************************************************************************************************
        for pattern_cat_match in pattern_category:
            pattern_match_flag = '' 
            #*************************************
            #filtering based on pattern_category
            #*************************************
            pattern_filter_condition_df = pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$".format(pattern_cat_match),case=False)]           
            if not pattern_filter_condition_df.empty:
                #**********************************************************
                #filtering based on pattern_keys and filter_condition
                #**********************************************************
                and_condition_df = pattern_filter_condition_df[pattern_filter_condition_df['filter_condition'].str.contains("^\s*{}\s*$".format('1'),case=False)]
                #print(and_condition_df['pattern_keys'].values.tolist())
                #print('and_condition_df',and_condition_df)
                or_condition_df = pattern_filter_condition_df[pattern_filter_condition_df['filter_condition'].str.contains("^\s*{}\s*$".format('0'),case=False)]
                #*****************************'*********************************************
                #checking if both filter condition '1' or '0' present in the category
                #**************************************************************************
                if  all([and_condition_df.empty,or_condition_df.empty]) == True:                 
                    and_condition_list = and_condition_df['pattern_keys'].values.tolist()                                        
                    or_condition_list = or_condition_df['pattern_keys'].values.tolist()
                    
                    if all(match.lower().strip() in content.lower() for match in and_condition_list):
                        if any(match.lower() in content.lower() for match in or_condition_list):
                            pattern_match_flag = 's'
                           # print('kamal',pattern_cat_match)
                #******************************************
                #checking only filter condition 'and' 
                #******************************************

                elif not and_condition_df.empty:
                    and_condition_list = and_condition_df['pattern_keys'].values.tolist()
                    if all(match.lower().strip() in content.lower() for match in and_condition_list):
                        pattern_match_flag = 's'
                #******************************************
                #checking only filter condition 'or' 
                #******************************************
                elif not or_condition_df.empty:
                    or_condition_list = or_condition_df['pattern_keys'].values.tolist()
                    #print(or_condition_list)
                    if any(match.lower().strip() in content.lower() for match in or_condition_list):
                        pattern_match_flag = 's'
            #****************************************************
            #if the file under goes any of the pattern category
            #****************************************************                             
            if  pattern_match_flag == 's':      
                logger.info('{} is found in {}'.format(file,pattern_cat_match))
                #*****************************************************************************************
                #validating for relavent and non-relavant file if the key-value extract is text from file
                #*****************************************************************************************
                    
                if all(int(match.lower().strip()) == 0 for match in (pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$".format(pattern_cat_match),case=False)])['result_type'].values.tolist()):
                    sil_elast_product_list = None
                    for match_mpm_cat in silicone_elastomer_product_df.columns:
                        if match_mpm_cat.lower().strip() == pattern_cat_match.strip().lower():                            
                            sil_elast_product_list = list(set(silicone_elastomer_product_df[match_mpm_cat].values.tolist())) 
                    file_loc = file.rsplit('all-text',1)[0] + 'valid-files/'+ pattern_cat_match.strip() + '/'
                    file_is_valid = file_is_valid_query.format(1,1,'null',file.replace("dbfs:","/dbfs"))
                    update_operation(file_is_valid,sql_conn,cursor)
                    relavent_text_extract(file,file_loc,content,product_inscope_df,pattern_cat_match.strip(),file_is_valid_query,file_unique_list,sql_conn,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query,sil_elast_product_list)       
                    #file_type_list.append('text')
                    file_valid_flag ='s'

                #******************************************************************************************
                #validating for relavent and non-relavant file if the key-value extract is image from file
                #******************************************************************************************  
                elif all(int(match.lower().strip()) == 1 for match in (pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$".format(pattern_cat_match),case=False)])['result_type'].values.tolist()):                  
                    file_loc = file.rsplit('all-text',1)[0] + 'valid-files/'+ pattern_cat_match.strip() + '/'
                    
                    #*********************************************************************************
                    #file: will hold the file present in the statging path for process using tesseract
                    #*********************************************************************************
                    print(file.replace("dbfs:","/dbfs"))
                    file_is_valid = file_is_valid_query.format(1,1,'null',file.replace("dbfs:","/dbfs"))
                    update_operation(file_is_valid,sql_conn,cursor)
                    file = external_staging_path[index] + file.rsplit('/',1)[1].rsplit('.',1)[0] + external_satging_file_format[index]
                    #print(file) 
                    relavent_image_extract(file,file_loc,content,product_inscope_df,pattern_cat_match.strip(),file_is_valid_query,file_unique_list,sql_conn,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query)
                    file_valid_flag ='s'

                #******************************************************************************************
                #validating for relavent and non-relavant file if the key-value extract  from excel
                #******************************************************************************************  
                elif all(int(match.lower().strip()) == 2 for match in (pattern_key_df[pattern_key_df['pattern_category'].str.contains("^\s*{}\s*$".format(pattern_cat_match),case=False)])['result_type'].values.tolist()):
                    head, tail = os.path.split(file)
                    file_extn = tail.rsplit('.',1)[-1]
                    file_name = tail.rsplit('.',1)[0]
                    file_in_dir =  os.listdir(file.rsplit('all-text',1)[0])
                    file_loc = file.rsplit('all-text',1)[0] + 'valid-files/'+ pattern_cat_match.strip() + '/'
                    excel_valid_query = config.get('mount_path','excel_is_valid')
                  
                    if file_name + '.csv' in file_in_dir:            
                        file_csv = file.rsplit('all-text',1)[0] + file_name + '.csv'
                        file_csv = file_csv.replace("/dbfs","dbfs:")
                        file_loc = file_loc.replace("/dbfs","dbfs:")   
                        dbutils.fs.cp(file_csv, file_loc) 
                        file_valid_path = file_loc + file_name + '.csv'
                        excel_valid_query.format(1,file_valid_path, pattern_cat_match.strip() ,file)
                        #file_type_list.append('records')
                    elif file_name in  file_in_dir:
                        file_csv_list = glob.glob(file.rsplit('all-text',1)[0] + file_name + '/*.csv')
                        for file_csv in file_csv_list:
                          file_name = file_csv.rsplit('.',1)[1]
                          file_csv = file_csv.replace("/dbfs","dbfs:")
                          file_loc = file_loc.replace("/dbfs","dbfs:")   
                          dbutils.fs.cp(file_csv, file_loc)
                          file_valid_path = file_loc + file_name
                          excel_valid_query.format(1,file_valid_path, pattern_cat_match.strip(),file) 
                          #file_type_list.append('records')
                    file_valid_flag ='s'   
                  
        #*************************************************************************************************
        #Moving the files to invalid-files folder as file content doesn't fall under pattern validation
        #*************************************************************************************************
        if file_valid_flag !='s':
            file_loc = file.rsplit('all-text',1)[0] + 'invalid-files/'   
            
            if file_loc.strip() not in invalid_folder_list:
              path_exists(file_loc)
              invalid_folder_list.append(file_loc.strip())
            file = file.replace("/dbfs","dbfs:")
            file_loc = file_loc.replace("/dbfs","dbfs:") 
            dbutils.fs.cp(file, file_loc) 
            file_name = file.rsplit('/',1)[-1]
            file_loc = file_loc.replace("dbfs:","/dbfs") + file_name
            file_is_valid = file_is_valid_query.format(0,0,file_loc,file.replace("dbfs:","/dbfs"))
            update_operation(file_is_valid,sql_conn,cursor)  

      except Exception as e:
        logger.error('file not found {}'.format(file),exc_info=True)
    unstruct_data_df['file_path']  = file_path_list 
    unstruct_data_df['file_type']  = file_type_list
    unstruct_data_df['category']  = category_list 
    unstruct_data_df['product_type']  = product_type_list
    unstruct_data_df['product']  = product_valid_list 
    return unstruct_data_df
  except Exception as e:
    logger.error('something went wrong in pattern match validation',exc_info=True)
    
    
#**************************************************************************************************************
#function name: excel2csv
#Ojective: excel to csv formats
#input parameters
#path: will hold the xlsx file path 
#Sheet: will hold sheet name present in the xlsx file
#Usage:converts excel file type into csv for text extarction as excel not supported in databricks
#called by: xlsx_text_extract
#**************************************************************************************************************
def excel2csv(path, sheet):
  try:
      wb = openpyxl.load_workbook(path)
      sh = wb[sheet]
      head, tail = os.path.split(path)
      filename = path.split('/')[-1].split('.')[0]
      file = head + '/' + 'temp/csv/' + sheet + '.csv'
      dbutils.fs.mkdirs((absolute_path +'temp/csv/').replace("/dbfs","dbfs:")) 
      with open(file, 'w', encoding="utf-8") as f:
          c = csv.writer(f)
          for r in sh.rows:
              c.writerow([cell.value for cell in r])
      return file
  except Exception as e:
    logger.error('Error in excel2csv function while converting {}'.format(path),exc_info=True)
    
#********************************************************************************************************************
#Function name: excel2txt
#Objective: Excel to text 
#input parameters
#staging_path : will hold the staging path of EXCEl type file fetched 
#abs_path: will hold file to be extracted 
#filename: will hold name of the file to be extracted
#Usage: common code is written which converts all the excel type files into text will be done using this function
#called by: csv_text_extract, xlsx_text_extract
#*********************************************************************************************************************
def excel2txt(staging_path, abs_path, filename, sheet):
  try:
    data_text = pd.read_csv(abs_path, encoding='cp1252')
    file = staging_path +'temp/temp_all_text/'+ filename +'_'+ sheet+'.txt'
    dbutils.fs.mkdirs((staging_path +'temp/temp_all_text/').replace("/dbfs","dbfs:")) 
    data_text.to_csv(file)
    
  except Exception as e:
     logger.error('Error in excel2csv function while converting {}'.format(abs_path),exc_info=True)

#******************************************************************************************************************************
#function name: csv_text_extract
#Ojective: csv file into text extract
#input parameters
#staging_path : will hold the staging path of xlsx file fetched 
#csv_list: will hold all the csv files in a list
#source_type: will hold the respective source type of staging path
#all_files: will hold all-text path where text files to be stored which got extracted from the csv sheet
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#Usage: common code is written which converts all the csv file into text and stores the extracted data in all_files area in txt format,  #       then file path into the file_processing_info table
#called by : external_folder_structure_process
#****************************************************************************************************************************** 
def csv_text_extract(staging_path,csv_list,source_type,all_files,file_processing_info,sql_conn,cursor):
  try:
      logger.info('Executing csv_text_extract function') 
      for abs_path in csv_list:
        try:
          head, tail = os.path.split(abs_path)
          file_extn = cstail.rsplit('.',1)[-1]
          file_name = tail.rsplit('.',1)[0]
          dbutils.fs.rm((staging_path +'temp/temp_all_text/').replace("/dbfs",""),True)
          dbutils.fs.rm((staging_path +'temp/csv/').replace("/dbfs",""),True)
          sheet =str(1)
          #**********************************************************************************************************************
          #excel2txt: It converts CSV file into text by taking  file path and sheet name present in the csv as input  
          #**********************************************************************************************************************
          excel2txt(staging_path, abs_path, file_name, sheet)
          text = glob.glob(staging_path +'temp/temp_all_text/'+'*.txt')
          dbutils.fs.mkdirs(all_files.replace("/dbfs","dbfs:")) 
          text_csv = pd.DataFrame()
          file_path = all_files + file_name + '.txt'
          for t in text:
            data = pd.read_csv(t, encoding='utf-8')
            text_csv = text_csv.append(data)
          text_csv.to_csv(file_path)
            #*******************************************************************************************************************
            #Creation of insert query for the extracted valid file path to the file_processing_info table and executed using
            #update_operation
            #*******************************************************************************************************************
          file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, {})".format(source_type, file_name, 
            'Excel','.csv', staging_path.replace('//','/'), file_path.replace('//','/'), 1,0,'GETDATE()','GETDATE()')
          update_operation(file_processing_info_query,sql_conn,cursor)
          logger.error('{}  extract_csv_text sucessfully'.format(file_path.replace('//','/')))
            
        except Exception as e:
          #********************************************************************************************************************
          #Creation of insert query for the extracted invalid file path to the file_processing_info table and executed using
          #update_operation
          #********************************************************************************************************************
          file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, {})".format(source_type, file_name, 
              'Excel','.csv', staging_path.replace('//','/'), 'null', 0,0,'null','null')
          update_operation(file_processing_info_query,sql_conn,cursor)
          logger.error('Error in extracting csv_ text {}'.format(file_path.replace('//','/')))
          
  except Exception as e:
    logger.error('Something went wrong in the csv_text_extract function', exc_info=True)    
#*****************************************************************************************************************************************
#function name : xlsx_text_extract
#objective : Extraction of text from excel sheets
#file_processing_info table
#input parameters
#staging_path : will hold the staging path of xlsx file fetched 
#xlsx_list: will hold all the xlsx files in a list
#source_type: will hold the respective source type of staging path
#all_files: will hold all-text path where text files to be stored which got extracted from the excel sheet
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database
#Usage: common code is written which converts all the xlsx file into text and stores the extracted data in all_files area in txt format,  #       then file path into the file_processing_info table
#called by : external_folder_structure_process
#*****************************************************************************************************************************************      
def xlsx_text_extract(staging_path,xlsx_list,source_type,all_files,file_processing_info,sql_conn,cursor):
  try:
      for abs_path in xlsx_list:
        try:
          head, tail = os.path.split(abs_path)
          file_extn = tail.rsplit('.',1)[-1]
          file_name = tail.rsplit('.',1)[0]
          wb = openpyxl.load_workbook(abs_path) 
          allsheets = list(wb.sheetnames)
          dbutils.fs.rm((staging_path +'temp/temp_all_text/').replace("/dbfs",""),True)
          dbutils.fs.rm((staging_path +'temp/csv/').replace("/dbfs",""),True)
          file_path = all_files + file_name + '.txt'
          for sheet in allsheets:
            excel2csv(abs_path, sheet)
          temp_path = glob.glob(staging_path+'temp/csv/'+'*.*')
          dbutils.fs.mkdirs((all_files + file_name +'/').replace("/dbfs","dbfs:"))
          dbutils.fs.cp((staging_path +'temp/csv/').replace("/dbfs","dbfs:"), (all_files + file_name+'/').replace("/dbfs","dbfs:"), recurse=True)
          for i in range(len(temp_path)):
              excel2txt(staging_path, temp_path[i], file_name, sheet)
              text_excel = glob.glob(staging_path +'temp/temp_all_text/'+'*.txt')
              dbutils.fs.mkdirs(all_files.replace("/dbfs","dbfs:"))
              text1 = pd.DataFrame()
              for t in text_excel:
                data = pd.read_csv(t, encoding='utf-8')
                text1 = text1.append(data)
              text1.to_csv(file_path)
              
          #**************************************************************************************************************
          #Creation of insert query for the extracted valid file path to the file_processing_info table using
          #update_operation
          #*************************************************************************************************************** 
          file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, {})".format(source_type, file_name,
          'Excel','.csv', staging_path.replace('//','/'), file_path.replace('//','/'), 1,0,'GETDATE()','GETDATE()')
          update_operation(file_processing_info_query,sql_conn,cursor)
          logger.error('{}  extract_csv_text sucessfully'.format(file_path.replace('//','/')))          
                    
        except Exception as e:     
            #**************************************************************************************************************
            #Creation of insert query for the extracted invalid file path to the file_processing_info table using
            #update_operation
            #***************************************************************************************************************
            file_processing_info_query = file_processing_info + " values ('{}', '{}', '{}', '{}', '{}', '{}', '{}', '{}', {}, {})".format(source_type, file_name, 
              'Excel','.csv', staging_path.replace('//','/'), 'null', 0,0,'null','null')
            update_operation(file_processing_info_query,sql_conn,cursor)
            logger.error('{}  is not extracted'.format(file_path.replace('//','/')))
            logger.error('Error in xlsx_text_extract function while converting {}'.format(abs_path),exc_info=True)
          
  except Exception as e:
    logger.error('Something went wrong in the xlsx_text_extract function', exc_info=True)       
    
#**************************************************************************************************************************************
#Function name: external_folder_structure_process
#objective: To convert all the file types into text format 
#input Parameters:
#external_folder_structure: Will call external_source_data function by passing external_folder_structure_query and returns all the
#                           details in the external_folder_structure in a dataframe
#external_source_file_formats: Will call external_source_data function by passing file_format_query and returns all the
#                              inscope file formats to be consider for processing
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#sql_conn: will hold DB_connectivity 
#Cursor: will hold cursor object for executing queries, it helps to update the database                
#Usage: Extraction of input file data fetched from the external_folder_structure table and stores all the extracted file path in the
#file_processing_info table
#ouput: returns raw_df which holds all the staging file path in dataframe which helps to move file to processed folder
#**************************************************************************************************************************************
def external_folder_structure_process(external_folder_structure,external_source_file_formats,file_processing_info,sql_conn,cursor):
  try:
    #raw_df will be used for moving the raw files into processed folder after the key-data extract
    #raw files and raw format will conatin each raw file path and formats in the list
    #after all the file extraction raw files and raw format will append to the raw_df dataframe
    raw_df1 = pd.DataFrame()
    #raw_df.columns = ['file_name']
    raw_files = []
    raw_format = []  
    print(external_folder_structure)
    for index in external_folder_structure.index:
      source_type = external_folder_structure['blob_src_type'][index].strip()
      mount_path = external_folder_structure['db_fs_mount_path'][index].strip()
      staging_path =  mount_path + external_folder_structure['absolute_path'][index]
      if os.path.exists(staging_path):
          all_files= staging_path.split('staging',1)[0] + 'analytics/processed/all-text/'  
          #***************************************************************************
          #fetching all the pdf file types from the sources
          #pdf_file_list: will have all the pdf file path from each category
          #***************************************************************************
          if '.pdf' in external_source_file_formats:
              pdf_file_list = glob.glob(staging_path +'*.pdf')       
              if bool(pdf_file_list):
                  logger.info('{} pdf files found in the {}'.format(len(pdf_file_list),staging_path))
                  raw_files = raw_files + pdf_file_list
                  raw_format = raw_format + ['.pdf']*len(pdf_file_list)
                  #***************************************************************************************************
                  #sharepoint_native_scanned_pdf_split: will split the pdf files into two types like(native, scanned)
                  #***************************************************************************************************
                  native_path, scanned_path = sharepoint_native_scanned_pdf_split(staging_path,pdf_file_list[15:17]) 
                  if native_path != None:
                   #***************************************************************************************************
                   #native_pdf_extract_text: will extract data from the native pdf type
                   #**************************************************************************************************
                      native_pdf_extract_text(native_path,all_files,staging_path,source_type,file_processing_info,sql_conn,cursor)
                  if scanned_path != None:
                   #***************************************************************************************************
                   #scanned_pdf_extract_text: will extract data from the scanned pdf type
                   #***************************************************************************************************
                      scanned_pdf_extract_text(scanned_path,all_files,staging_path,source_type,file_processing_info,sql_conn,cursor)
        #**************************************************************
        #fetching all the Document file types from the sources
        #**************************************************************        
          if '.docx' in  external_source_file_formats:
              doc_file_list = glob.glob(staging_path+'*.docx')
              if bool(doc_file_list):
                  logger.info('{} docx files found in the {}'.format(len(doc_file_list),staging_path))
                  raw_files = raw_files + doc_file_list
                  raw_format = raw_format + ['.docx']*len(doc_file_list)
                  #***************************************************************************************************
                  #extract_doc_text: will extract data from the documnet file type
                  #doc_file_list: will have all the document file path from each category
                  #***************************************************************************************************
                  extract_doc_text(staging_path,doc_file_list,source_type,all_files,file_processing_info,sql_conn,cursor)

        #******************************************************************************
        #fetching all the message file types from the sources
        #msg_list: will have all the message file path from each category
        #*******************************************************************************      
          if '.msg' in  external_source_file_formats:
              msg_list = glob.glob(staging_path+'*.msg')
              staging_path_index = staging_path.lower().find('staging')
              if staging_path_index != -1:
                staging_path_pdf = staging_path[:staging_path_index] + 'staging/pdf/raw/'
              if bool(msg_list):  
                 #*********************************************************************
                 #outlook_attachment: will fetch the attachments found in the messasge
                 #*********************************************************************              
                 pdf_out_look = outlook_attachment(msg_list,staging_path_pdf,raw_files,raw_format)
                 if bool(pdf_out_look):
                   native_path, scanned_path = sharepoint_native_scanned_pdf_split(staging_path_pdf,pdf_out_look) 
                   if native_path != None:
                     native_pdf_extract_text(native_path,all_files,staging_path_pdf,source_type,file_processing_info,sql_conn,cursor)
                   if scanned_path != None:
                     scanned_pdf_extract_text(scanned_path,all_files,staging_path_pdf,source_type,file_processing_info,sql_conn,cursor) 



        #**************************************************************
        #fetching all the csv file types from the sources
        #csv_list: will have all the csv format file path from each category
        #************************************************************** 
          if '.csv ' in external_source_file_formats:
              csv_list = glob.glob(staging_path+'*.csv')
              if bool(csv_list):
                #*********************************************************************
                #csv_text_extract: will extract the data from the csv file type
                #*********************************************************************
                logger.info('{} csv file found in the staging_path'.format(len(csv_list)))
                csv_text_extract(staging_path,csv_list,source_type,all_files,file_processing_info,sql_conn,cursor)
        #*******************************************************************************
        #fetching all the xlsx and xlsm file types from the sources
        #xlsx_list: will have all the xlsx format file path from each category
        #xlsm_list: will have all the xlsm format file path from each category
        #*******************************************************************************
          if 'xlsx' or 'xlsm' in external_source_file_formats:
              xlsx_list = glob.glob(staging_path+'*.xlsx')
              xlsm_list = glob.glob(staging_path+'*.xlsm')
              xlsx_list = xlsx_list + xlsm_list
              #****************************************************************************
              #xlsx_text_extract: will extract the data from the xlsx and xlsm file type
              #****************************************************************************
              if bool(xlsx_list):
                  logger.info('{} xlsx file found in the staging_path'.format(len(xlsx_list)))
                  xlsx_text_extract(staging_path,xlsx_list,source_type,all_files,file_processing_info,sql_conn,cursor)
    raw_df1['file_name'] = raw_files 
    return raw_df1
  except Exception as e:
    logger.error('Something went wrong in the external_folder_structure_process function {} ', exc_info=True)
#********************************************************************************************************************************
#function name: external_source_data
#Objective:Select operations on the table
#input Parameters: sql_conn will hold DB_connectivity object and query will hold select operations or query to fetch the data from 
#                  the table
#ouput Parameter: result will hold the fetched data from azure sql table in a dataframe
#called by: main, pattern_match_validation function
#Usage: common function is written to perform select query operation on the required table and returns the result in dataframe
#********************************************************************************************************************************          
def external_source_data(sql_conn,query):
  try:
    if sql_conn is not None:  
      result = pd.read_sql(query, sql_conn) 
      logger.info('Successfully extracted the data of momentive.external_source_folder_structure from sql server')      
    else:
      logger.error('Sql_conn has None value something went wrong in the Sql server connection') 
    
    return result
  except Exception as error:
    logger.error('Something went wrong in the external_source_data function', exc_info=True)

#*************************************************************************************************************************************
#function name: update_operation
#Objective: insert, update and Delete operations on the table
#Usage: common function is written to perform (insert, update and Delete) query operation on the required table 
#input Parameters: sql_conn will hold DB_connectivity object and Cursor will hold cursor object for executing queries, it helps to 
#                  update the database
#called by: native_pdf_extract_text, scanned_pdf_extract_text, extract_doc_text, xlsx_text_extract, csv_text_extract,file_validation #functions   
#*************************************************************************************************************************************
def update_operation(query,sql_conn,cursor):
 # print(query)
  cursor.execute(query)
  sql_conn.commit()  
  
#**********************************************************************************************
#function name: Sql_db_connection
#Objective: connecting sql db using pyodbc
#Usage: common function is written to connect with given database using pyodbc package
#output: Sql_conn will hold the DB_connectivity object
#called by : Main function
#**********************************************************************************************
def Sql_db_connection(): 
  try:
    server = config.get('sql_db', 'server')
    database = config.get('sql_db', 'database')
    username = config.get('sql_db', 'username')
    password = config.get('sql_db', 'password')
    DATABASE_CONFIG = {
      'server': server,
      'database': 'cld-it-dev-pih-db1',
      'username': username,
      'password': password
    }
    print(DATABASE_CONFIG)
    driver= "{ODBC Driver 17 for SQL Server}"
    connection_string = 'DRIVER=' + driver + \
                      ';SERVER=' + DATABASE_CONFIG['server'] + \
                      ';PORT=1433' + \
                      ';DATABASE=' + DATABASE_CONFIG['database'] + \
                      ';UID=' + DATABASE_CONFIG['username'] + \
                      ';PWD=' + DATABASE_CONFIG['password'] 


    sql_conn = pyodbc.connect(connection_string)
    #result = pd.read_sql('select * from  momentive.category_keywords', sql_conn)
    #print(result['field'])
    logger.info('Successfully connected with the sql serevr ')
    if sql_conn is None:
      logger.error('sql is not connected')    
    return sql_conn    
  except Exception as e:
    logger.error('Something went wrong in the Sql_db_connection function', exc_info=True)


#****************************************************************************************************************************************** 
#Function name: Main 
#Objective: Program will start process using this function 
#sql_conn: Azure SQl DB Connectivity will be created using this  Sql_db_connection()   
#cursor: Cursor will be created using this sql_conn.cursor() for executing Sql operations
#external_source_folder_structure table: external_source_folder_structure table will contain all the data ingestion details from azure to 
#                                  blob storgae based on each category  
#external_folder_structure_query: Will have the query to connect with  external_source_folder_structure table
#file_format table: will have all the inscope file formats(like 'PDF', "Document') etc to be processed
#file_format_query: inscope file format query will get capture in this variable
#file_processing_info table: will have all the extracted file path details, Sources type (like Sharepoint,Website etc) and  
#file_processing_info: Will have query for updating the processed file information in the file_processing_info table
#external_source_data: Will perform sql select operation by passing two parameter DB Connectivity(sql_conn) and select query.
#                      it will return output in dataframe 
#external_folder_structure: Will call external_source_data function by passing external_folder_structure_query and returns all the
#                           details in the external_folder_structure in a dataframe
#external_source_file_formats: Will call external_source_data function by passing file_format_query and returns all the
#                              inscope file formats to be consider for processing
#external_folder_structure_process: its a function which takes 5 parameters like(external_folder_structure,external_source_file_formats,
#                                   file_processing_info, ,sql_conn, sql_conn) using this will etract all the file data got from the  
#                                   external_folder_structure and stores the extracted path in the file_processing_info table
#external_file_process_query: 
#pattern_match_validation: it will do the pattern matching for 15 categories for extracted files using file information stored in the    
#                           file_processing_info table and pattern matching keywords stored in the pattern_matching_keys table
#*****************************************************************************************************************************************  
   
def main(): 
  try:
      sql_conn = Sql_db_connection()
      cursor = sql_conn.cursor()  
      external_folder_structure_query = config.get('mount_path', 'external_source_folder_structure')
      file_format_query = config.get('mount_path', 'external_source_file_formats')
      file_processing_info = config.get('mount_path', 'file_processing_info')
      external_folder_structure = external_source_data(sql_conn,external_folder_structure_query)
      external_source_file_formats = external_source_data(sql_conn,file_format_query)['file_format'].values.tolist()
      raw_df = external_folder_structure_process(external_folder_structure,external_source_file_formats,file_processing_info,sql_conn,cursor)    
      external_file_process_query = config.get('mount_path', 'external_file_process')
      external_processed_files_df = external_source_data(sql_conn,external_file_process_query)
      unstruct_category_key_query = config.get('mount_path','unstruct_category_key')
      unstruct_category_key_df = external_source_data(sql_conn,unstruct_category_key_query)
      unstructure_processed_data_query = config.get('mount_path', 'unstructure_processed_data')
      #unstructure_processed_data(unstructure_processed_data_query,unstruct_key_value_data_df,sql_conn,cursor)
      pattern_match_validation(sql_conn,external_processed_files_df,cursor,unstruct_category_key_df,raw_df,unstructure_processed_data_query)
      #unstruct_key_value_data_df = key_value_extract(unstruct_data_key_info,raw_df,unstruct_category_key_df)
      #unstruct_key_value_data_df.to_csv(('/dbfs/mnt/momentive-sources-pih/sharepoint-pih/customer-communications-pih/mpm-2019-pih/') +'metadata.csv')
      #unstructure_processed_data_query = config.get('mount_path', 'unstructure_processed_data')
      #unstructure_processed_data(unstructure_processed_data_query,unstruct_key_value_data_df,sql_conn,cursor)
  except Exception as e:
    logger.error('Something went wrong in main function', exc_info=True)
    
  
#***************************************************************************************** 
#calling the main function when python code is triggered from azure pipeline                                      
#***************************************************************************************** 

if __name__ == '__main__':
  logger.info('Main function Started')
  main()

# COMMAND ----------

/dbfs/mnt/momentive-sources-pih/sharepoint-pih/customer-communications-pih/mpm-2019-pih/analytics/processed/all-text/Chemical Structures.txt

# COMMAND ----------

# MAGIC %sh
# MAGIC cat /databricks/driver/shared_main_code_error.log

# COMMAND ----------

